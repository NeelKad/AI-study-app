from werkzeug.security import generate_password_hash, check_password_hash
from flask_sqlalchemy import SQLAlchemy
from datetime import datetime, timedelta
from functools import wraps
from werkzeug.utils import secure_filename
import PyPDF2
import docx
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
import youtube_dl
import tempfile

# SSL fix
os.environ['SSL_CERT_FILE'] = certifi.where()

app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY", "supersecretkey")

# File upload config
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'doc', 'pptx', 'ppt'}
MAX_FILE_SIZE = 16 * 1024 * 1024  # 16MB

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE

# Create upload directory
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Database config
DATABASE_URL = os.getenv("DATABASE_URL")
if DATABASE_URL and DATABASE_URL.startswith("postgres://"):
    DATABASE_URL = DATABASE_URL.replace("postgres://", "postgresql://", 1)
app.config['SQLALCHEMY_DATABASE_URI'] = DATABASE_URL or 'sqlite:///site.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
db = SQLAlchemy(app)

# OpenAI client
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY", ""))

login_manager = LoginManager()
login_manager.login_view = "login"
login_manager.init_app(app)

# --- MODELS ---
class User(UserMixin, db.Model):
    __tablename__ = 'users'
    id = db.Column(db.Integer, primary_key=True)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password = db.Column(db.String(256), nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    trial_expires_at = db.Column(db.DateTime, nullable=True)
    has_unlimited_access = db.Column(db.Boolean, default=False)

    def get_id(self):
        return str(self.id)
    
    def is_trial_expired(self):
        if self.has_unlimited_access:
            return False
        if self.trial_expires_at is None:
            return True
        return datetime.utcnow() > self.trial_expires_at
    
    def get_time_remaining(self):
        if self.has_unlimited_access:
            return "Unlimited"
        if self.trial_expires_at is None:
            return "No trial expiry set"
        if self.is_trial_expired():
            return "Expired"
        remaining = self.trial_expires_at - datetime.utcnow()
        minutes = int(remaining.total_seconds() // 60)
        seconds = int(remaining.total_seconds() % 60)
        return f"{minutes}m {seconds}s"

class Subject(db.Model):
    __tablename__ = 'subjects'
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('users.id'), nullable=False)
    name = db.Column(db.String(255), nullable=False)
    color = db.Column(db.String(7), default='#6366f1')  # Hex color
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    
    user = db.relationship('User', backref='subjects')
    notes = db.relationship('Note', backref='subject', cascade='all, delete-orphan')

class Note(db.Model):
    __tablename__ = 'notes'
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('users.id'))
    subject_id = db.Column(db.Integer, db.ForeignKey('subjects.id'), nullable=True)
    title = db.Column(db.String(255))
    content = db.Column(db.Text)
    source_type = db.Column(db.String(50), default='manual')  # manual, pdf, docx, pptx, website, youtube
    source_url = db.Column(db.String(500), nullable=True)  # For websites/youtube
    file_path = db.Column(db.String(500), nullable=True)  # For uploaded files
    created_at = db.Column(db.DateTime, server_default=db.func.now())
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)
    
    user = db.relationship('User', backref='notes')

class StudyScheduleItem(db.Model):
    __tablename__ = 'study_schedule_items'
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('users.id'))
    title = db.Column(db.String(255), nullable=False)
    description = db.Column(db.Text)
    study_action = db.Column(db.String(100))  # flashcards, questions, summarise, etc.
    note_id = db.Column(db.Integer, db.ForeignKey('notes.id'), nullable=True)
    priority = db.Column(db.String(20), default='medium')  # low, medium, high
    estimated_time = db.Column(db.Integer, default=15)  # minutes
    completed = db.Column(db.Boolean, default=False)
    is_ai_generated = db.Column(db.Boolean, default=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    scheduled_for = db.Column(db.DateTime, nullable=True)
    
    user = db.relationship('User', backref='schedule_items')
    note = db.relationship('Note', backref='schedule_items')

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

@app.before_request
def create_tables():
    db.create_all()

# --- HELPER FUNCTIONS ---
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(file_path):
    """Extract text from PDF file"""
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text()
        return text
    except Exception as e:
        print(f"Error extracting PDF text: {e}")
        return ""

def extract_text_from_docx(file_path):
    """Extract text from Word document"""
    try:
        doc = docx.Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting DOCX text: {e}")
        return ""

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint presentation"""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting PPTX text: {e}")
        return ""

def extract_text_from_website(url):
    """Extract text from website"""
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Get text and clean it up
        text = soup.get_text()
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = ' '.join(chunk for chunk in chunks if chunk)
        
        return text[:5000]  # Limit to 5000 characters
    except Exception as e:
        print(f"Error extracting website text: {e}")
        return ""

def extract_text_from_youtube(url):
    """Extract transcript from YouTube video"""
    try:
        ydl_opts = {
            'writesubtitles': True,
            'writeautomaticsub': True,
            'subtitleslangs': ['en'],
            'skip_download': True,
        }
        
        with tempfile.TemporaryDirectory() as temp_dir:
            ydl_opts['outtmpl'] = os.path.join(temp_dir, '%(title)s.%(ext)s')
            
            with youtube_dl.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(url, download=False)
                title = info.get('title', 'YouTube Video')
                
                # Try to get subtitles
                subtitles = info.get('subtitles', {})
                auto_captions = info.get('automatic_captions', {})
                
                text = f"Title: {title}\n\n"
                
                # Extract subtitle text if available
                if 'en' in subtitles or 'en' in auto_captions:
                    # This is a simplified approach - in reality, you'd need to download and parse subtitle files
                    text += "Transcript not available through this method. Please use YouTube's built-in transcript feature."
                else:
                    text += "No transcript available for this video."
                
                return text
    except Exception as e:
        print(f"Error extracting YouTube text: {e}")
        return f"Error processing YouTube video: {str(e)}"

# --- DECORATOR ---
def trial_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not current_user.is_authenticated:
            return redirect(url_for('login'))
        if current_user.is_trial_expired():
            flash('Your trial has expired! Enter admin key to unlock.', 'error')
            return redirect(url_for('trial_expired'))
        return f(*args, **kwargs)
    return decorated_function

# --- SUBJECT ROUTES ---
@app.route('/api/subjects', methods=['GET'])
@login_required
@trial_required
def api_get_subjects():
    subjects = Subject.query.filter_by(user_id=current_user.id).all()
    return jsonify([{
        'id': s.id,
        'name': s.name,
        'color': s.color,
        'note_count': len(s.notes)
    } for s in subjects])

@app.route('/api/subjects', methods=['POST'])
@login_required
@trial_required
def api_create_subject():
    data = request.json
    subject = Subject(
        user_id=current_user.id,
        name=data.get('name'),
        color=data.get('color', '#6366f1')
    )
    db.session.add(subject)
    db.session.commit()
    return jsonify({'id': subject.id, 'name': subject.name, 'color': subject.color})

@app.route('/api/subjects/<int:subject_id>', methods=['PUT'])
@login_required
@trial_required
def api_update_subject(subject_id):
    subject = Subject.query.filter_by(id=subject_id, user_id=current_user.id).first()
    if not subject:
        return jsonify({'error': 'Subject not found'}), 404
    
    data = request.json
    subject.name = data.get('name', subject.name)
    subject.color = data.get('color', subject.color)
    db.session.commit()
    return jsonify({'message': 'Subject updated'})

@app.route('/api/subjects/<int:subject_id>', methods=['DELETE'])
@login_required
@trial_required
def api_delete_subject(subject_id):
    subject = Subject.query.filter_by(id=subject_id, user_id=current_user.id).first()
    if not subject:
        return jsonify({'error': 'Subject not found'}), 404
    
    db.session.delete(subject)
    db.session.commit()
    return jsonify({'message': 'Subject deleted'})

# --- FILE UPLOAD ROUTES ---
@app.route('/api/upload', methods=['POST'])
@login_required
@trial_required
def api_upload_file():
    try:
        data = request.form
        title = data.get('title', '').strip()
        subject_id = data.get('subject_id') or None
        source_type = data.get('source_type', 'manual')
        
        content = ""
        source_url = None
        file_path = None
        
        if source_type == 'file' and 'file' in request.files:
            file = request.files['file']
            if file and file.filename and allowed_file(file.filename):
                filename = secure_filename(file.filename)
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], 
                                       f"{current_user.id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{filename}")
                file.save(file_path)
                
                # Extract text based on file type
                ext = filename.rsplit('.', 1)[1].lower()
                if ext == 'pdf':
                    content = extract_text_from_pdf(file_path)
                    source_type = 'pdf'
                elif ext in ['docx', 'doc']:
                    content = extract_text_from_docx(file_path)
                    source_type = 'docx'
                elif ext in ['pptx', 'ppt']:
                    content = extract_text_from_pptx(file_path)
                    source_type = 'pptx'
                
                if not title:
                    title = filename.rsplit('.', 1)[0]
        
        elif source_type == 'website':
            url = data.get('url', '').strip()
            if url:
                content = extract_text_from_website(url)
                source_url = url
                if not title:
                    title = f"Website: {url[:50]}..."
        
        elif source_type == 'youtube':
            url = data.get('url', '').strip()
            if url:
                content = extract_text_from_youtube(url)
                source_url = url
                if not title:
                    title = f"YouTube: {url[:50]}..."
        
        elif source_type == 'manual':
            content = data.get('content', '').strip()
        
        if not content.strip():
            return jsonify({'error': 'No content could be extracted'}), 400
        
        if not title:
            title = f"Untitled {source_type.title()}"
        
        # Create note
        note = Note(
            user_id=current_user.id,
            subject_id=int(subject_id) if subject_id else None,
            title=title,
            content=content,
            source_type=source_type,
            source_url=source_url,
            file_path=file_path
        )
        
        db.session.add(note)
        db.session.commit()
        
        return jsonify({
            'message': 'Note created successfully',
            'note_id': note.id,
            'title': note.title
        })
    
    except Exception as e:
        db.session.rollback()
        if 'file_path' in locals() and file_path and os.path.exists(file_path):
            os.remove(file_path)
        return jsonify({'error': str(e)}), 500

# --- NOTE ROUTES ---
@app.route('/api/notes/<int:note_id>', methods=['PUT'])
@login_required
@trial_required
def api_update_note(note_id):
    note = Note.query.filter_by(id=note_id, user_id=current_user.id).first()
    if not note:
        return jsonify({'error': 'Note not found'}), 404
    
    data = request.json
    note.title = data.get('title', note.title)
    note.content = data.get('content', note.content)
    note.subject_id = data.get('subject_id') or None
    note.updated_at = datetime.utcnow()
    
    db.session.commit()
    return jsonify({'message': 'Note updated successfully'})

# --- AI SCHEDULE GENERATION ---
def generate_ai_study_schedule(user_notes):
    """Generate AI-powered study schedule based on user's notes"""
    if not user_notes:
        return []
    
    # Create a summary of all notes for the AI
    notes_summary = ""
    for note in user_notes:
        notes_summary += f"Title: {note.title}\nContent: {note.content[:200]}...\n\n"
    
    try:
        prompt = f"""
        Based on the following study notes, create a personalized study schedule. Generate 5-8 study tasks that will help the student learn effectively.

        Notes:
        {notes_summary}

        For each study task, provide:
        1. A clear title
        2. A brief description of what to study
        3. Recommended study action (flashcards, questions, summarise, tutor, or pastpaper)
        4. Priority level (high, medium, low)
        5. Estimated time in minutes (10-60)

        Prioritize recent notes and complex topics. Mix different study methods for variety.
        
        Respond with a JSON array where each item has these fields:
        - title: string
        - description: string
        - study_action: string (flashcards/questions/summarise/tutor/pastpaper)
        - priority: string (high/medium/low)
        - estimated_time: integer (minutes)
        - note_title: string (which note this relates to)
        """

        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=1000
        )
        
        ai_response = response.choices[0].message.content.strip()
        # Extract JSON from response (in case there's extra text)
        start_idx = ai_response.find('[')
        end_idx = ai_response.rfind(']') + 1
        if start_idx != -1 and end_idx != -1:
            json_str = ai_response[start_idx:end_idx]
            schedule_items = json.loads(json_str)
            return schedule_items
        else:
            return []
    except Exception as e:
        print(f"Error generating AI schedule: {e}")
        return []

# --- SCHEDULE API ROUTES ---
@app.route('/api/generate-schedule', methods=['POST'])
@login_required
@trial_required
def api_generate_schedule():
    """Generate AI study schedule"""
    try:
        # Get user's notes
        user_notes = Note.query.filter_by(user_id=current_user.id).all()
        
        if not user_notes:
            return jsonify({"error": "No notes found to generate schedule from"}), 400
        
        # Generate AI schedule
        ai_schedule = generate_ai_study_schedule(user_notes)
        
        if not ai_schedule:
            return jsonify({"error": "Failed to generate schedule"}), 500
        
        # Clear existing AI-generated items
        StudyScheduleItem.query.filter_by(user_id=current_user.id, is_ai_generated=True).delete()
        
        # Save new schedule items to database
        saved_items = []
        for item in ai_schedule:
            # Find matching note
            note = None
            if 'note_title' in item:
                note = Note.query.filter_by(
                    user_id=current_user.id, 
                    title=item['note_title']
                ).first()
                if not note:
                    # Try partial match
                    note = Note.query.filter(
                        Note.user_id == current_user.id,
                        Note.title.like(f"%{item['note_title']}%")
                    ).first()
            
            schedule_item = StudyScheduleItem(
                user_id=current_user.id,
                title=item.get('title', 'Study Task'),
                description=item.get('description', ''),
                study_action=item.get('study_action', 'summarise'),
                priority=item.get('priority', 'medium'),
                estimated_time=item.get('estimated_time', 15),
                note_id=note.id if note else None,
                is_ai_generated=True
            )
            db.session.add(schedule_item)
            saved_items.append({
                'title': schedule_item.title,
                'description': schedule_item.description,
                'study_action': schedule_item.study_action,
                'priority': schedule_item.priority,
                'estimated_time': schedule_item.estimated_time
            })
        
        db.session.commit()
        return jsonify({"message": "Schedule generated successfully", "items": saved_items})
        
    except Exception as e:
        db.session.rollback()
        return jsonify({"error": str(e)}), 500

@app.route('/api/add-schedule-item', methods=['POST'])
@login_required
@trial_required
def api_add_schedule_item():
    """Add custom schedule item"""
    data = request.json
    
    try:
        schedule_item = StudyScheduleItem(
            user_id=current_user.id,
            title=data.get('title', ''),
            description=data.get('description', ''),
            study_action=data.get('study_action', 'summarise'),
            priority=data.get('priority', 'medium'),
            estimated_time=int(data.get('estimated_time', 15)),
            note_id=data.get('note_id') if data.get('note_id') else None,
            is_ai_generated=False
        )
        
        db.session.add(schedule_item)
        db.session.commit()
        
        return jsonify({"message": "Schedule item added successfully"})
    except Exception as e:
        db.session.rollback()
        return jsonify({"error": str(e)}), 500

@app.route('/api/toggle-schedule-item/<int:item_id>', methods=['POST'])
@login_required
@trial_required
def api_toggle_schedule_item(item_id):
    """Toggle completion status of schedule item"""
    try:
        item = StudyScheduleItem.query.filter_by(
            id=item_id, 
            user_id=current_user.id
        ).first()
        
        if not item:
            return jsonify({"error": "Schedule item not found"}), 404
        
        item.completed = not item.completed
        db.session.commit()
        
        return jsonify({"message": "Schedule item updated", "completed": item.completed})
    except Exception as e:
        db.session.rollback()
        return jsonify({"error": str(e)}), 500

@app.route('/api/delete-schedule-item/<int:item_id>', methods=['DELETE'])
@login_required
@trial_required
def api_delete_schedule_item(item_id):
    """Delete schedule item"""
    try:
        item = StudyScheduleItem.query.filter_by(
            id=item_id, 
            user_id=current_user.id
        ).first()
        
        if not item:
            return jsonify({"error": "Schedule item not found"}), 404
        
        db.session.delete(item)
        db.session.commit()
        
        return jsonify({"message": "Schedule item deleted"})
    except Exception as e:
        db.session.rollback()
        return jsonify({"error": str(e)}), 500

# --- AUTH ROUTES ---
@app.route('/signup', methods=['GET', 'POST'])
def signup():
    if request.method == 'POST':
        email = request.form['email']
        password = request.form['password']
        if User.query.filter_by(email=email).first():
            flash('Email already exists.', 'error')
            return render_template('signup.html')
        hashed_pw = generate_password_hash(password)
        trial_expiry = datetime.utcnow() + timedelta(minutes=10)  # 10 min trial
        new_user = User(email=email, password=hashed_pw, trial_expires_at=trial_expiry)
        db.session.add(new_user)
        db.session.commit()
        flash('Account created! Please log in.', 'success')
        return redirect(url_for('login'))
    return render_template('signup.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        email = request.form['email']
        password = request.form['password']
        user = User.query.filter_by(email=email).first()
        if user and check_password_hash(user.password, password):
            if user.trial_expires_at is None and not user.has_unlimited_access:
                user.trial_expires_at = datetime.utcnow() + timedelta(minutes=10)
                db.session.commit()
            login_user(user)
            return redirect(url_for('dashboard'))
        flash('Invalid credentials.', 'error')
    return render_template('login.html')

@app.route('/logout')
@login_required
def logout():
    logout_user()
    return redirect(url_for('login'))

# --- TRIAL EXPIRED ---
@app.route('/trial-expired')
@login_required
def trial_expired():
    return render_template('trial_expired.html', user=current_user)

@app.route('/trial-expired/unlock', methods=['POST'])
@login_required
def trial_unlock():
    ADMIN_KEY = os.getenv("ADMIN_KEY", "your-secret-admin-key-here")
    entered_key = request.form.get('admin_key', '')

    if entered_key != ADMIN_KEY:
        flash('Invalid admin key!', 'error')
        return redirect(url_for('trial_expired'))

    user = current_user
    user.has_unlimited_access = True
    db.session.commit()

    flash('Unlimited access granted! Welcome back.', 'success')
    return redirect(url_for('dashboard'))

# --- ADMIN ROUTES ---
@app.route('/admin/grant-access', methods=['GET', 'POST'])
def admin_grant_access():
    ADMIN_KEY = os.getenv("ADMIN_KEY", "your-secret-admin-key-here")
    if request.method == 'POST':
        provided_key = request.form.get('admin_key')
        user_email = request.form.get('user_email')
        if provided_key != ADMIN_KEY:
            flash('Invalid admin key!', 'error')
            return render_template('admin_grant_access.html')
        user = User.query.filter_by(email=user_email).first()
        if not user:
            flash('User not found!', 'error')
            return render_template('admin_grant_access.html')
        user.has_unlimited_access = True
        db.session.commit()
        flash(f'Unlimited access granted to {user_email}!', 'success')
        return render_template('admin_grant_access.html')
    return render_template('admin_grant_access.html')

@app.route('/admin/users')
def admin_users():
    ADMIN_KEY = os.getenv("ADMIN_KEY", "your-secret-admin-key-here")
    provided_key = request.args.get('key')
    if provided_key != ADMIN_KEY:
        return "Access denied", 403
    users = User.query.all()
    return render_template('admin_users.html', users=users)

# --- TRIAL STATUS API ---
@app.route('/api/trial-status')
@login_required
def api_trial_status():
    if current_user.has_unlimited_access:
        return jsonify({"unlimited": True, "expired": False, "time_remaining": "Unlimited", "remaining_seconds": -1})
    if current_user.is_trial_expired():
        return jsonify({"unlimited": False, "expired": True, "time_remaining": "Expired", "remaining_seconds": 0})
    if not current_user.trial_expires_at:
        return jsonify({"unlimited": False, "expired": True, "time_remaining": "No trial expiry set", "remaining_seconds": 0})
    remaining = current_user.trial_expires_at - datetime.utcnow()
    remaining_seconds = int(remaining.total_seconds())
    minutes = remaining_seconds // 60
    seconds = remaining_seconds % 60
    return jsonify({"unlimited": False, "expired": False, "time_remaining": f"{minutes}m {seconds}s", "remaining_seconds": remaining_seconds})

# --- DASHBOARD & NOTES ---
@app.route('/dashboard')
@login_required
@trial_required
def dashboard():
    notes = Note.query.filter_by(user_id=current_user.id).all()
    schedule_items = StudyScheduleItem.query.filter_by(user_id=current_user.id).order_by(
        StudyScheduleItem.completed.asc(),
        StudyScheduleItem.priority.desc(),
        StudyScheduleItem.created_at.desc()
    ).all()
    return render_template('dashboard.html', notes=notes, schedule_items=schedule_items, email=current_user.email, user=current_user)

@app.route('/enter-notes')
@login_required
@trial_required
def enter_notes():
    return render_template('index.html', user=current_user)

@app.route('/add-note', methods=['POST'])
@login_required
@trial_required
def add_note():
    title = request.form.get('title')
    content = request.form.get('content')
    if not title or not content:
        flash("Please provide both title and content.", "error")
        return redirect(url_for('enter_notes'))
    note = Note(user_id=current_user.id, title=title, content=content)
    db.session.add(note)
    db.session.commit()
    flash("Note saved successfully!", "success")
    return redirect(url_for('dashboard'))

@app.route('/delete-note/<int:note_id>', methods=['DELETE'])
@login_required
@trial_required
def delete_note(note_id):
    """Delete a note"""
    try:
        note = Note.query.filter_by(id=note_id, user_id=current_user.id).first()
        
        if not note:
            return jsonify({"error": "Note not found or access denied"}), 404
        
        # Also delete any schedule items related to this note
        StudyScheduleItem.query.filter_by(note_id=note_id, user_id=current_user.id).delete()
        
        db.session.delete(note)
        db.session.commit()
        
        return jsonify({"message": "Note deleted successfully"})
    except Exception as e:
        db.session.rollback()
        return jsonify({"error": str(e)}), 500

@app.route('/note/<int:note_id>')
@login_required
@trial_required
def view_note(note_id):
    note = Note.query.filter_by(id=note_id, user_id=current_user.id).first()
    if note:
        return render_template('view_note.html', title=note.title, content=note.content, note_id=note.id)
    flash("Note not found or access denied.", "error")
    return redirect(url_for('dashboard'))

@app.route('/')
def index():
    return redirect(url_for('dashboard')) if current_user.is_authenticated else redirect(url_for('login'))

@app.route('/save-note', methods=['POST'])
@login_required
@trial_required
def save_note_alias():
    return add_note()

# --- STUDY TOOLS (UI) ---
@app.route('/flashcards')
@app.route('/flashcards/<int:note_id>')
@login_required
@trial_required
def flashcards(note_id=None):
    note_content = None
    if note_id:
        note = Note.query.filter_by(id=note_id, user_id=current_user.id).first()
        if note:
            note_content = note.content
    return render_template('flashcards.html', note_content=note_content, user=current_user)

@app.route('/questions')
@app.route('/questions/<int:note_id>')
@login_required
@trial_required
def questions(note_id=None):
    note_content = None
    if note_id:
        note = Note.query.filter_by(id=note_id, user_id=current_user.id).first()
        if note:
            note_content = note.content
    return render_template('questions.html', note_content=note_content, user=current_user)

@app.route('/summarise')
@app.route('/summarise/<int:note_id>')
@login_required
@trial_required
def summarise(note_id=None):
    note_content = None
    if note_id:
        note = Note.query.filter_by(id=note_id, user_id=current_user.id).first()
        if note:
            note_content = note.content
    return render_template('summarise.html', note_content=note_content, user=current_user)

@app.route('/pastpaper')
@app.route('/pastpaper/<int:note_id>')
@login_required
@trial_required
def pastpaper(note_id=None):
    note_content = None
    if note_id:
        note = Note.query.filter_by(id=note_id, user_id=current_user.id).first()
        if note:
            note_content = note.content
    return render_template('pastpaper.html', note_content=note_content, user=current_user)

@app.route('/tutor')
@app.route('/tutor/<int:note_id>')
@login_required
@trial_required
def tutor(note_id=None):
    note_content = None
    if note_id:
        note = Note.query.filter_by(id=note_id, user_id=current_user.id).first()
        if note:
            note_content = note.content
    return render_template('tutor.html', note_content=note_content, user=current_user)

@app.route('/my-notes')
@login_required
@trial_required
def my_notes():
    notes = Note.query.filter_by(user_id=current_user.id).all()
    return render_template('my_notes.html', notes=notes)

# --- OPENAI API ROUTES ---
@app.route('/api/flashcards', methods=['POST'])
@login_required
@trial_required
def api_flashcards():
    notes = request.json.get('notes', '')
    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "user", "content": "Generate flashcards from these notes. Output a valid JSON object where keys are terms and values are definitions."},
                {"role": "user", "content": notes}
            ],
            temperature=0.7,
            max_tokens=700
        )
        text = response.choices[0].message.content.strip()
        flashcards = json.loads(text) if text else {}
    except Exception:
        flashcards = {}
    return jsonify(flashcards)

@app.route('/api/questions', methods=['POST'])
@login_required
@trial_required
def api_questions():
    notes = request.json.get('notes', '')
    prompt = f"Generate 20 concise questions from these notes:\n\n{notes}. Also, ensure that it only covers the key points of the notes. Also, make them short answer questions rather than single number or word answers."
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7,
        max_tokens=500
    )
    text = response.choices[0].message.content.strip()
    questions = [q.strip() for q in text.split('\n') if q.strip()]
    return jsonify(questions)

@app.route('/api/grade_question', methods=['POST'])
@login_required
@trial_required
def api_grade_question():
    data = request.json
    question = data.get('question', '')
    answer = data.get('answer', '')
    prompt = (
        f"Grade this answer: '{answer}' for the question: '{question}'. "
        "Reply EXACTLY in this format:\n"
        "score: <number from 0 to 10>\n"
        "improvement: <suggestion>\n"
        "model answer: <model answer>"
    )
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7,
        max_tokens=200
    )
    text = response.choices[0].message.content.strip()
    score_match = re.search(r'score:\s*([0-9]+(?:\.[0-9]+)?)', text, re.IGNORECASE)
    improvement_match = re.search(r'improvement:\s*(.*?)(?:\nmodel answer:|$)', text, re.IGNORECASE | re.DOTALL)
    model_answer_match = re.search(r'model answer:\s*(.*)', text, re.IGNORECASE | re.DOTALL)
    return jsonify({
        "grade_score": score_match.group(1) if score_match else "N/A",
        "improvement": improvement_match.group(1).strip() if improvement_match else "No suggestion.",
        "model_answer": model_answer_match.group(1).strip() if model_answer_match else "No model answer provided."
    })

@app.route('/api/summarise', methods=['POST'])
@login_required
@trial_required
def api_summarise():
    notes = request.json.get('notes', '')
    prompt = f"Summarize these notes into a concise yet elaborate paragraph"
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7,
        max_tokens=400
    )
    summary = response.choices[0].message.content.strip()
    return jsonify({"summary": summary})

@app.route("/api/pastpaper", methods=["POST"])
@login_required
@trial_required
def api_pastpaper():
    data = request.get_json()
    notes = data.get("notes", "")
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "user", "content": f"Generate a past paper from these notes:\n{notes}. If it is math or science related, provide 20 mcq questions, 30 calcualtion/shortanswer questions , and 2 extremely challenging problems. If it is something like Humanities or English, provide 20 MCQs, 10 short answers, and 1 extended response at the end. At the end, provide an answer key. Also, make the entire paper look good, it should look almost identical to a typical exam with space for writing answers and working out. It should have sections formated for clarity"}
        ]
    )
    past_paper = response.choices[0].message.content
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    for line in past_paper.split("\n"):
        pdf.multi_cell(0, 10, line)
    pdf_bytes = pdf.output(dest='S').encode('latin1')
    return send_file(BytesIO(pdf_bytes), mimetype="application/pdf", as_attachment=True, download_name="past_paper.pdf")


from flask import session

@app.route("/api/tutor_chat", methods=["POST"])
@login_required
@trial_required
def tutor_chat():
    message = request.json.get("message", "").strip()
    note_content = request.json.get("note_content", "").strip()

    if not message:
        return jsonify({"error": "Empty message"}), 400

    if "tutor_chat_history" not in session:
        session["tutor_chat_history"] = [
            {"role": "system", "content": "You are a helpful AI study tutor. Always explain clearly."},
            {"role": "system", "content": f"Here are the user's study notes:\n{note_content}"}
        ]

    session["tutor_chat_history"].append({"role": "user", "content": message})

    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=session["tutor_chat_history"],
            temperature=0.7,
            max_tokens=500
        )
        reply = response.choices[0].message.content.strip()
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    session["tutor_chat_history"].append({"role": "assistant", "content": reply})
    session.modified = True

    return jsonify({"reply": reply})


if __name__ == '__main__':
    app.run(debug=True)







